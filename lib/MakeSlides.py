#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Wed Aug  7 15:06:54 2024

@author: dac
"""

## -*- coding: utf-8-*g
__doc__='''
#==============================================================================
# lib_wiskdac.py
@ last: 2024-03-16-2011z
@ author: Daniel.Collins <Daniel.Collins@wisk.aero >
# library of python tools shared for Wisk applications
#Python ; #wisk; #Daniel.COllins; #makeSlides
#------------------------------------------------------------------------------'''
###############################################################################
PSH =[0 for i in range(0,100)] ; ## PSH[1]  = 0 ;

PATH_SLIDEMASTER = r'/Users/dac/pylab/lib/';

##############################################################################
# In[0]: Load packages
#------------------------------------------------------------------------------
####  Import Packages
# import lib_tcpro as tcpro ;
# import lib_wiskdac as wisk ;
# import lib_getfiles as getfiles ;
# import tcpro_tempstat as tempstat;
# import showmaxia as showmaxia  ;
####  Import Packages
import re ;
from   os import listdir;
import os;
import string ;
import math ;
import linecache ;
import argparse ;
import statistics ;
import pandas as pd ;
import numpy as np ;
import matplotlib as mpl ;
import matplotlib.pyplot as plt ;
import matplotlib.font_manager as fontmgr ;
from   scipy import signal;
from   datetime import datetime ;
from   pptx import Presentation ;
from   PIL import Image ; ### what is the python -m pip install [] for PIL ??
#-----------------------------------------------------------------------------#

###############################################################################
def getfilesbe ( inppath, begkw='', endkw='.csv', PSH=[0, 0 ] ):
    '''
    Return FILESETPATHS <list of strings>
    Parameters
    ----------
    inppath [str]
        DESCRIPTION.
    begkw [str]
        keyword for files to be sought in path, def = None
    endkw [str]
        keyword for end of file name, def='csv'
    PSH [list]
        toggles for PrintOut
        *default is [0, 0 ].

    '''
    ### from: eslib_getfiles.py as getfiles.wisk.getfilebegins
    FILESETPATHS = [];
    if PSH[1] : print(f' inppath = {inppath}');
    for f in listdir(inppath):
      if f.endswith( endkw ) and f.startswith( begkw ) :
         # for ii in range(10,1,-1):
         #   ol = ''
         #   for jj in range(ii):
         #     ol += (r'\\')
         #   f = f.replace( ol , '');
           # eval( f' f = f.replace( (r'\\')*ii, '') ')
         FILESETPATHS.append( str( inppath + f ) );
    if PSH[1] : print( FILESETPATHS );
    return FILESETPATHS;
#-----------------------------------------------------------------------------#
###############################################################################
def getsimcase( inp, method='log' ): ## requires python v > 3.10
    if method=='rfile':
            simfile = os.path.split(inp)[1];
            simfile = simfile.rstrip( 'run.log' );
    if method=='sim':
            simfile = getfilesbe( inp, endkw='sim' );
    if method=='log': 
            simfile = getfilesbe( inp, endkw='run.log' );
            if type(simfile)==list:
                simfile = simfile[-1];
            simfile = simfile.split(inp)[1].rstrip( 'run.log' ) ;

    simfile = simfile.rstrip( 'run.log' );
    simfile = simfile.rstrip( '.' );
    simfile = simfile.rstrip( '-' );
    simfile = simfile.rstrip( '~' );
    simfile = simfile.rstrip( '.sim' );
    return simfile ;

#-----------------------------------------------------------------------------#
# In[1]: ##>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
def MakeSlides( ##>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
    ppt_gallery, 
    #EXE, 
    PSH,
    inppath=os.getcwd(),
    ncol=2,
    nrow=2,
    basename=None,
    xtagname=None,
    figb=4.000, ## base of pw in inches
    figr=None,  ## ph / pw ratio
    figp=1.333, ## ratio of pWidth / pHeight
    ph_set=None, ## manually define each frame's height
    pw_set=None, ## manually define each frame's width
    
    pwhmax=False, ## if True, make maximum width and height
    phmax=False, ## if True, make maximum width and height
    pwmax=True, ## if True, make maximum width and height

    PPTXPATH=PATH_SLIDEMASTER,
    PPTXFILE=r"SLIDES.pptx",
    SLIDE_TITLE = 1,
    SLIDE_BLANK = 1,    
    googledrive = None, # r"G:\My Drive\\" ,
    pptsavename = None,

    
    ):
    ''' 
SLIDENAMES =\
   MakeSlides.MakeSlides( ppt_gallery, inppath,, EXE, PSH, EXE, PSH, ...*optionals )

 » SLIDENAMES<list>: append (pptfile) file name of generated pptx
 *« ppt_gallery<list>: list of figures/png aka photo gallery
~~ optional:
  « inppath<str>: path of where input is occuring (..evnt None=> list.curdir) 
  « PSH<list>: toggles for print-outs
  « ncol=2 <int>: nr of columns framed on a slide
  « nrow=2 <int>: nr of rows   framed on a slide
  « figr=None <float>: ratio of height/width of the figures
  « figb=4.000 <float>: base of pw in inches
  « basename=None <str>: 
  « xtagname=None <str>: 
  « PPTXPATH=PYTHONPATH <str>: path locating the master pptxfile
  « PPTXFILE=r"SLIDES.pptx" <str>: file serving as Master Template
  « googledrive = r"G:\My Drive\\" <str>: path, cloned copy of generated Slide

? « EXE<dict>: toggles for execution
 Vers: dac/2024-02-16-1714Z~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'''
 
    print( ' Making Slide ....');
    from pptx      import Presentation ;
    from pptx.util import Inches       ;
    from pptx.util import Pt           ;
    # def Slidegen ( imgix, outpath ): #    
    # SLIDENAMES = MakeSlides( ppt_gallery, EXE, PSH, inppath, FSS, )      
    
    SLIDE_TITLE = 1 ;
    # SLIDE_BLANK = 1,    

    if 1: ### sizing of frames and figures on layout ----------------------------------------------
        #ncol = 3 ;
        #nrow = 2 ;
        nlay  = int(ncol * nrow) ;
        pa    = float( 0.050 ) ;
        pt    = float( 0.150 ) ; 
        mg    = float( 0.025 ) ;
        # pw = float( 2.25 ) ; ph = float( 2.25 ) ; 

        if figp: 
            figr = ( 1/ figp );
 
        if figr == None :
          if figb == None: 
            figb = 2.250 ; 
          figr = 1.00/ 1.00 ;
        elif figr == 9:
          figb = 0.500     ; figr = (9.00/16.00 );
        elif figr == 3:
          figb = 0.81*2.00 ; figr = (3.00/ 4.00 );
        elif figr == 2 :
          figb = 1.00      ; figr = (2.00/ 1.00 );
        elif figr == 1.5 :
          figb = figb      ; figr = (1.50/ 1.00 );
        elif figr == 2 :
          figb = 1.00      ; figr = (2.00/ 1.00 );
        else:
          figb = 1.00*figb ; figr = (1.00/ 1.00 )*figr ;
        
       
        if 1:
          pw = float(( figb  * 1.00  ) / ncol) ;
          ph = float(( figb  * figr  ) / nrow) ;
        
        if pw_set:
          pw = pw_set ;  
        if ph_set:
          ph = ph_set ;  

        if Inches( nrow*ph ) >  7.250  : ph =  5.00/nrow ;
        if Inches( ncol*pw ) > 10.500  : pw = 10.00/ncol ;
        if pwhmax:
            ph =  7.50/nrow ;
            pw = 13.00/ncol ;
            
        if phmax:
            ph =  7.50/nrow ;
        if pwmax:
            pw = 13.00/ncol ;
            
        pt =0.5*(  7.500 - (nrow)*(ph + 0*mg ) )
        pa =0.5*( 13.333 - (ncol)*(pw + 0*mg ) )
            
#### --------------------------------------------------------------------------
    if inppath ==None: inppath=os.getcwd();
    if type(inppath)==list: inppath=str(inppath[0]); 
    print( "inppath : {}".format(inppath));
    FSS = os.sep;
    SLIDENAMES = [];
#    if PPTXPATH==None:  PPTXPATH  = r"G:\My Drivve\mylab\lib\Wisk\\" ;    
#    if PPTXFILE==None:  PPTXFILE  = "WISK.pptx" ;
    PPTMASTER = PPTXPATH + FSS +  PPTXFILE ; 
    print( '\n    pptx slide master: {} \n'.format(PPTMASTER))
    #_save_ppt = inppath + FSS ; 
    _keypics  = '' ; # inppath['leg'] +FSS;
    TCLOCPIC = [
    '''
    from current Presentation Slide, export picture of each TCLOC_KEY_png that represents each GRP#
    ''', ## Grp.0
    'image of TC LOC in Group 1' ,
    'image of TC LOC in Group 2' ,
    '... etc ...',
    ];
    if PSH[0]: print(' now doing EXEPPT ');

    if (1==1):
        prs     = Presentation( PPTMASTER );
        # SLIDE_TITLE = 0,
        # SLIDE_BLANK = 7,    
        title_slide_layout = prs.slide_layouts[ SLIDE_TITLE ];
        blank_slide_layout = prs.slide_layouts[ SLIDE_BLANK ];
        # prs.slides.add_slide(blank_slide_layout) ;

        for i, g in enumerate(  ppt_gallery ):
          pic = g ;
          fo  = int (  i % nlay     );
          fr  = int (( i /ncol)%nrow);
          fc  = int ( i % ncol      );
          if PSH[3]: print( i, fo, fr, fc      );
            
          if 1: 
            if PSH[1]: print ('{} {}'.format(i, g ));
            # grp = int(g[4]);
            # if ( i % 2 )==0:
            # slide = prs.slides.add_slide(blank_slide_layout);
            if PSH[2]: print( ' Picture being processed: {}'.format( g ));
            # if ( i % pframelayout ) < pframelayout:
            #     pass;
            if ( i % nlay ) == 0 :
               slide = prs.slides.add_slide(blank_slide_layout);
               shapes = slide.shapes ;
              # shapes.title.text = '{:s}'.format( str(g[9]));
            if 1: 
                #( int ( ra / ncol) == 0 ) :
              slide.shapes.add_picture(
                g,#g[0],
                Inches( pa +fc*pw ), Inches(pt + fr* ph),
                width=Inches(pw), height=Inches(ph),
                ) ;
            '''
                left = top = width = height = Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = "This is text inside a textbox"
            ###-------------------------------------------------      '''
        ###EO for ------------------------------------------
        ### cfg with multiple groups ... 
        ###  if ( i==len(ppt_gallery)-1 or (i == limGrp-1)):
        
        if 1: ## naming generated pptxfile 
            import shutil; 
            if basename==None : 
              try:
                  basename = getsimcase( inppath );
              except:
                  basename = '__Slides__';
            if xtagname==None: xtagname='';
            if basename==None: basename=''; 
            if pptsavename==None:
                pptfile = \
                + 1 * basename \
                + 1 * ( '_' +  xtagname )\
                + 1 * '.pptx' ;
            else: pptfile  = pptsavename + '.pptx' ;
            print( 'inppath :', inppath );
            print( 'pptfile :', pptfile );
            print( 'pptsavename :', pptsavename );
            pptfilefull = inppath + FSS + pptfile ;
            #if PSH[1]: 
            print(f' saving file {pptfilefull}');
 
            if 1:
                prs.save( pptfilefull );
                SLIDENAMES.append( pptfile );
            #googledrive = r"G:\My Drive\Review\\" ;
            if 0:
                if googledrive==None: 
                    pass; 
                else:
                    gcopy = googledrive + FSS + pptfile ;
                    shutil.copy2( pptfilefull , str(googledrive + FSS + pptfile) );
                    print(' pptx made here and {}'.format(  gcopy ));
            # print( os.listdir( inppath ) ) ; 
    ###########################################################################
    return SLIDENAMES  ;
####<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<####
####EOF MakeSlides ############################################################
helpinfo_Module_pptx='''
from pptx      import Presentation ;
from pptx.util import Inches       ;
from pptx.util import Pt           ;
# from pptx.text import text  ;
python -m pip install python-pptx
https://python-pptx.readthedocs.io/en/latest/index.html
https://python-pptx.readthedocs.io/en/latest/user/quickstart.html#add-picture-example
https://python-pptx.readthedocs.io/en/latest/user/quickstart.html#add-textbox-example
''';